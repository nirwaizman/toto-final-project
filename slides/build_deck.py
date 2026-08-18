"""Rebuild the final project slide deck with real charts, screenshots, and clean design."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DARK = RGBColor(0x0F, 0x17, 0x2A)
CARD = RGBColor(0x1E, 0x29, 0x3B)
BLUE = RGBColor(0x25, 0x63, 0xEB)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
GREEN = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xCB, 0xD5, 0xE1)

CH = "slides/charts"
SCREEN = "slides/screenshot-app.png"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def add_bg(slide, color=DARK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_text(slide, l, t, w, h, text, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font="Calibri", line_spacing=1.15):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return box


def add_bullets(slide, l, t, w, h, items, size=15, color=GRAY, bullet_color=ORANGE, bold_first=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"●  {item}"
        p.line_spacing = 1.3
        p.space_after = Pt(10)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = "Calibri"
    return box


def add_accent_bar(slide, color=ORANGE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), SH)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_header(slide, title, subtitle=None, accent=ORANGE):
    add_accent_bar(slide, accent)
    add_text(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.8),
              title, size=30, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.05), Inches(12), Inches(0.5),
                  subtitle, size=15, color=GRAY)


def add_pagenum(slide, n):
    add_text(slide, Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.35),
              str(n), size=11, color=GRAY, align=PP_ALIGN.RIGHT)


def img_card(slide, path, l, t, w, h, caption=None):
    """Place image inside a rounded card with border."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h + (Inches(0.35) if caption else 0))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
    card.line.width = Pt(1)
    card.shadow.inherit = False
    pad = Inches(0.1)
    slide.shapes.add_picture(path, l + pad, t + pad, height=h - 2*pad)
    if caption:
        add_text(slide, l, t + h - Inches(0.05), w, Inches(0.4), caption,
                  size=11, color=GRAY, align=PP_ALIGN.CENTER)


slide_num = 1

# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_accent_bar(s, BLUE)
add_text(s, Inches(1), Inches(2.3), Inches(11.3), Inches(1.4),
          "Toto — Industry-Simulated AI Product Workflow", size=40, bold=True, color=WHITE)
add_text(s, Inches(1), Inches(3.7), Inches(11.3), Inches(0.6),
          "A Two-Crew CrewAI System for Hotel Booking Cancellation Prediction", size=20, color=ORANGE)
add_text(s, Inches(1), Inches(4.5), Inches(11.3), Inches(0.6),
          "Final Project — AI Development & Collaboration Course", size=16, color=GRAY)
add_text(s, Inches(1), Inches(4.95), Inches(11.3), Inches(0.5),
          "Nir Waizman  |  August 2026", size=16, color=GRAY)
add_pagenum(s, slide_num); slide_num += 1

# ---------- Slide 2: Agenda ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Agenda")
agenda = [
    "Business Background & the Problem",
    "The Concept — meet Toto",
    "High-Level Architecture (2 Crews + Flow)",
    "Dataset & EDA",
    "Crew 1 — Data Analyst (4 agents)",
    "Crew 2 — Data Scientist (3 agents)",
    "Model Results & Model Card",
    "Streamlit App Demo",
    "Tech Stack & Conclusion",
]
add_bullets(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5), agenda, size=18)
add_pagenum(s, slide_num); slide_num += 1

# ---------- Slide 3: Business Background ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Business Background")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(6.2), Inches(5.2), [
    "A hospitality group operates multiple hotels (inspired by real operations: boutique hotel in Eilat + city-center hotel in Tel Aviv).",
    "Booking cancellations create major revenue-management risk: unsold rooms, wasted staffing, inaccurate forecasts.",
    "The business needs a data-driven way to flag high-risk bookings early — before check-in — to enable smarter overbooking and deposit policy.",
    "Goal: build an end-to-end, reproducible AI pipeline that ingests booking data, extracts insights, and predicts cancellation risk.",
], size=16)
img_card(s, f"{CH}/class_balance.png", Inches(7.2), Inches(1.6), Inches(5.4), Inches(4.9),
          caption="119,390 real hotel bookings (Portugal, 2015–2017)")
add_pagenum(s, slide_num); slide_num += 1

# ---------- Slide 4: The Concept — Toto ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, 'The Concept — "Toto"', accent=GREEN)
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.7), Inches(5), [
    "Inspired by Toto: a real, daily-operating personal AI assistant already handling legal tracking, banking, family logistics and hotel operations for its user.",
    "This project simulates the same philosophy at industry scale: specialized AI agents, each with one clear job, collaborating through defined handoffs.",
    "Crew 1 (Data Analyst) mirrors how Toto ingests raw, messy real-world information and turns it into structured, trustworthy insight.",
    "Crew 2 (Data Scientist) mirrors how Toto turns insight into a decision-support tool — a live prediction, not just a report.",
    "The CrewAI Flow enforces validation gates between crews — just like Toto never acts on unverified data.",
], size=17)
add_pagenum(s, slide_num); slide_num += 1

# ---------- Slide 5: Architecture ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "High-Level Architecture")

def arch_box(l, t, w, h, title, lines, color):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = CARD
    box.line.color.rgb = color; box.line.width = Pt(2)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; p.text = title
    p.runs[0].font.size = Pt(16); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = color
    for line in lines:
        p2 = tf.add_paragraph(); p2.text = "• " + line
        p2.runs[0].font.size = Pt(12); p2.runs[0].font.color.rgb = GRAY
        p2.space_before = Pt(4)

arch_box(Inches(0.6), Inches(1.7), Inches(3.7), Inches(3.6), "Crew 1\nData Analyst",
          ["Ingestion", "Cleaning", "EDA", "Insights + Contract"], BLUE)
arch_box(Inches(4.85), Inches(2.4), Inches(3.6), Inches(2.2), "CrewAI Flow",
          ["Validation gate", "Fails gracefully on\nleakage / mismatch", "Auto-handoff"], ORANGE)
arch_box(Inches(9.05), Inches(1.7), Inches(3.7), Inches(3.6), "Crew 2\nData Scientist",
          ["Contract Validator", "Feature Engineer", "ML Engineer (train + eval)"], GREEN)

# arrows
for x1, x2, y in [(Inches(4.3), Inches(4.85), Inches(3.5)), (Inches(8.45), Inches(9.05), Inches(3.5))]:
    conn = s.shapes.add_connector(2, x1, y, x2, y)
    conn.line.color.rgb = WHITE
    conn.line.width = Pt(2.5)

add_text(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.2),
          "Output: clean_data.csv → eda_report.html → insights.md → dataset_contract.json → features.csv → model.pkl → evaluation_report.md → model_card.md",
          size=13, color=GRAY)
add_pagenum(s, slide_num); slide_num += 1

# ---------- Slide 6: Dataset ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Dataset")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(5.8), Inches(5), [
    "Hotel Booking Demand — 119,390 real hotel reservations.",
    "Resort Hotel + City Hotel, Portugal, 2015–2017.",
    "Temporal split: train on 2015–2016 (77,438 rows), test on 2017 (40,129 rows) — simulates real forecasting.",
    "Leakage columns explicitly excluded: reservation_status, reservation_status_date.",
    "33 columns in the clean data; 25 contract-approved modeling features.",
], size=16)
img_card(s, f"{CH}/bookings_by_month.png", Inches(6.7), Inches(1.6), Inches(5.9), Inches(5.1),
          caption="Seasonal booking volume — summer peak clearly visible")
add_pagenum(s, slide_num); slide_num += 1

# ---------- Slide 7: Crew 1 ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Crew 1 — Data Analyst (4 Agents)")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(6.0), Inches(5), [
    "Data Ingestion Specialist — profiles the raw dataset, flags types and ranges.",
    "Data Cleaning Specialist — handles nulls, duplicates, outliers.",
    "EDA Analyst — generates statistics and an HTML report.",
    "Insights & Contract Writer — writes business insights and produces the machine-readable dataset_contract.json used to gate Crew 2.",
], size=16)
img_card(s, f"{CH}/cancel_by_hotel.png", Inches(7.0), Inches(1.6), Inches(5.6), Inches(5.1),
          caption="City Hotel cancels far more often than Resort Hotel")
add_pagenum(s, slide_num); slide_num += 1

# ---------- Slide 8: Key Insights ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Key Business Insights", accent=GREEN)
add_bullets(s, Inches(0.7), Inches(1.6), Inches(5.9), Inches(5.2), [
    "City Hotel cancels 41.8% of bookings vs 28.1% for Resort Hotel — nearly 50% higher risk.",
    "Non-refundable deposits show a paradoxically high cancellation rate — a red flag worth investigating operationally.",
    "Longer lead time correlates strongly with higher cancellation probability.",
    "Recommendation: apply stricter overbooking buffers for City Hotel and long-lead-time bookings.",
], size=16)
img_card(s, f"{CH}/cancel_by_deposit.png", Inches(6.9), Inches(1.6), Inches(5.7), Inches(5.1),
          caption="Cancellation rate by deposit type")
add_pagenum(s, slide_num); slide_num += 1

# ---------- Slide 9: Flow handoff ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "CrewAI Flow — The Handoff", accent=ORANGE)
add_bullets(s, Inches(0.7), Inches(1.6), Inches(6.0), Inches(5), [
    "Automates Crew 1 → Crew 2 handoff — no manual file passing.",
    "Validation gate: confirms clean_data.csv matches dataset_contract.json before any modeling starts.",
    "Fails gracefully: if leakage columns are detected or the contract doesn't match, the Flow stops Crew 2 and reports why.",
    "This mirrors production ML pipelines, where a broken upstream contract must never silently corrupt downstream models.",
], size=16)
img_card(s, f"{CH}/lead_time_dist.png", Inches(7.0), Inches(1.6), Inches(5.6), Inches(5.1),
          caption="Lead time distribution: canceled vs completed bookings")
add_pagenum(s, slide_num); slide_num += 1

# ---------- Slide 10: Crew 2 ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Crew 2 — Data Scientist (3 Agents)", accent=GREEN)
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.8), Inches(4.8), [
    "Contract Validator — confirms data matches the contract, flags any leakage columns before proceeding.",
    "Feature Engineer — builds the modeling feature set (features.csv) from validated clean data.",
    "ML Engineer — trains and evaluates candidate models, selects the best one, writes evaluation_report.md and model_card.md.",
], size=18)
add_pagenum(s, slide_num); slide_num += 1

# ---------- Slide 11: Model Results ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Model Results")
img_card(s, f"{CH}/model_comparison.png", Inches(0.7), Inches(1.55), Inches(8.2), Inches(5.2))
add_bullets(s, Inches(9.15), Inches(1.7), Inches(3.5), Inches(5), [
    "Compared RandomForest vs GradientBoosting.",
    "Temporal train/test split (2015–16 → 2017).",
    "GradientBoosting selected: ROC-AUC 0.875, F1 0.693.",
    "Edges out RandomForest on ROC-AUC (0.875 vs 0.871).",
], size=14)
add_pagenum(s, slide_num); slide_num += 1

# ---------- Slide 12: Model Card ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Model Card — Responsible AI", accent=ORANGE)
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.8), Inches(5), [
    "Purpose: support overbooking / deposit-policy decisions — not to deny individual guests service.",
    "Model: GradientBoostingClassifier, trained on 2015–2016 bookings, tested on 2017 (temporal holdout).",
    "Performance: Accuracy 78.7%, Precision 79.1%, Recall 61.7%, F1 0.693, ROC-AUC 0.875.",
    "Known limitation: recall of ~62% means ~38% of true cancellations are missed — model should support, not replace, human judgment.",
    "No PII used as a feature; country and market segment are aggregate categorical fields only.",
], size=17)
add_pagenum(s, slide_num); slide_num += 1

# ---------- Slide 13: Streamlit ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Streamlit Application", accent=GREEN)
add_bullets(s, Inches(0.7), Inches(1.55), Inches(4.4), Inches(5.2), [
    "Overview tab — explains the architecture.",
    "Run Flow tab — triggers the live CrewAI Flow end-to-end.",
    "EDA tab — interactive charts on the cleaned dataset.",
    "Model tab — live cancellation-risk prediction from user input.",
], size=15)
img_card(s, SCREEN, Inches(5.3), Inches(1.55), Inches(7.3), Inches(5.3))
add_pagenum(s, slide_num); slide_num += 1

# ---------- Slide 14: Tech Stack ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Tech Stack & Reproducibility")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.8), Inches(5), [
    "CrewAI, Python 3.12, GitHub, Streamlit.",
    "Pandas, Scikit-Learn, Matplotlib/Seaborn for data + modeling.",
    "Claude (Anthropic) as the LLM backing every agent.",
    "Fully reproducible: requirements.txt, .env-based secrets (excluded from Git), documented outputs contract.",
    "Public repo: github.com/nirwaizman/toto-final-project",
], size=17)
add_pagenum(s, slide_num); slide_num += 1

# ---------- Slide 15: Conclusion ----------
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Conclusion", accent=BLUE)
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.8), Inches(5), [
    "A working, reproducible two-crew CrewAI Flow — 7 specialized agents total.",
    "Solves a real hospitality problem: predicting booking cancellation risk to support revenue decisions.",
    "Full pipeline: raw data → validated contract → engineered features → trained model → live app.",
    "Same philosophy as Toto: specialized agents, verified handoffs, decision support — not blind automation.",
], size=18)
add_pagenum(s, slide_num); slide_num += 1

os.makedirs("slides", exist_ok=True)
prs.save("slides/toto-final-project.pptx")
print("Saved", len(prs.slides.__iter__.__self__._sldIdLst), "slides" if False else f"{len(prs.slides._sldIdLst)} slides")
